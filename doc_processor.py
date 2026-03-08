#!/usr/bin/env python3
"""Ocean Eterna Document Preprocessor — converts documents to indexed text.

Supports: PDF, DOCX, XLSX, CSV, PPTX, images (OCR), Jupyter notebooks,
code files (.py, .r, .js, .cpp, .java, .sql, .html), TXT, MD.

Usage:
    python3 doc_processor.py report.pdf                     # ingest into OE
    python3 doc_processor.py *.docx                         # batch ingest
    python3 doc_processor.py --dry-run scan.png             # preview extraction
    python3 doc_processor.py --server http://host:9090 f.pdf  # custom server
"""
import os
import re
import csv
import sys
import json
import uuid
import shutil
import argparse
from pathlib import Path
from datetime import datetime, timezone

import requests

OE_BASE = os.environ.get("OE_BASE_URL", "http://localhost:9090")
MAX_CONTENT_SIZE = 10_000_000  # 10MB text cap

# all supported extensions
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".csv",
    ".pptx",
    ".png", ".jpg", ".jpeg",
    ".txt", ".md",
    ".py", ".r", ".js", ".ts", ".cpp", ".c", ".h", ".hpp",
    ".java", ".sql", ".go", ".rs", ".rb", ".sh", ".bash",
    ".ipynb",
    ".html", ".htm",
}

# extension -> category mapping
CATEGORY_MAP = {
    ".pdf": "document", ".docx": "document", ".txt": "document", ".md": "document",
    ".html": "document", ".htm": "document",
    ".xlsx": "spreadsheet", ".csv": "spreadsheet",
    ".pptx": "presentation",
    ".png": "image", ".jpg": "image", ".jpeg": "image",
    ".py": "code", ".r": "code", ".js": "code", ".ts": "code",
    ".cpp": "code", ".c": "code", ".h": "code", ".hpp": "code",
    ".java": "code", ".sql": "code", ".go": "code", ".rs": "code",
    ".rb": "code", ".sh": "code", ".bash": "code",
    ".ipynb": "code",
}

# extension -> language name (for code files)
LANGUAGE_MAP = {
    ".py": "python", ".r": "r", ".js": "javascript", ".ts": "typescript",
    ".cpp": "cpp", ".c": "c", ".h": "c-header", ".hpp": "cpp-header",
    ".java": "java", ".sql": "sql", ".go": "go", ".rs": "rust",
    ".rb": "ruby", ".sh": "bash", ".bash": "bash",
}


def normalize_paragraphs(text: str) -> str:
    """clean text to UTF-8 with consistent \\n\\n paragraph separators."""
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[^\S\n]+", " ", text)
    text = re.sub(r" +\n", "\n", text)
    return text.strip()


def get_category(ext: str) -> str:
    """return document category based on file extension."""
    return CATEGORY_MAP.get(ext, "document")


def generate_summary(content: str, filename: str, category: str) -> str:
    """generate a short summary from extracted content."""
    # strip markdown headings and whitespace for cleaner summary
    lines = []
    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        # skip page markers
        if stripped.startswith("[Page ") and stripped.endswith("]"):
            continue
        # remove markdown heading markers
        if stripped.startswith("#"):
            stripped = stripped.lstrip("#").strip()
        if stripped:
            lines.append(stripped)

    if not lines:
        return f"{filename} ({category})"

    # take first meaningful text, cap at 200 chars
    summary_text = " ".join(lines[:5])
    if len(summary_text) > 200:
        # cut at word boundary
        summary_text = summary_text[:197]
        last_space = summary_text.rfind(" ")
        if last_space > 150:
            summary_text = summary_text[:last_space]
        summary_text += "..."
    return summary_text


def generate_file_id() -> str:
    """generate a unique file ID."""
    return "f_" + uuid.uuid4().hex[:12]


# ── extractors ───────────────────────────────────────────────────────

def extract_pdf(filepath: str) -> str:
    try:
        import pymupdf
    except ImportError:
        raise ImportError("PDF support requires pymupdf: pip install pymupdf")

    try:
        doc = pymupdf.open(filepath)
    except Exception as e:
        raise ValueError(f"cannot open PDF '{filepath}': {e}")

    pages = []
    for page_num, page in enumerate(doc):
        text = page.get_text("text").strip()

        # scanned page — try OCR fallback
        if len(text) < 50:
            try:
                import pytesseract
                from PIL import Image

                pix = page.get_pixmap(dpi=300)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text = pytesseract.image_to_string(img, lang="eng").strip()
                pix = None
            except Exception:
                if not text:
                    text = f"[Page {page_num + 1}: could not extract text]"

        if text:
            pages.append(f"[Page {page_num + 1}]\n\n{text}")

    doc.close()
    return "\n\n".join(pages)


def extract_docx(filepath: str) -> str:
    try:
        from docx import Document
    except ImportError:
        raise ImportError("DOCX support requires python-docx: pip install python-docx")

    try:
        doc = Document(filepath)
    except Exception as e:
        raise ValueError(f"cannot open DOCX '{filepath}': {e}")

    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            level = para.style.name.replace("Heading ", "").strip()
            prefix = "#" * int(level) if level.isdigit() else "##"
            parts.append(f"{prefix} {text}")
        else:
            parts.append(text)

    # extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append("\n".join(rows))

    return "\n\n".join(parts)


def extract_xlsx(filepath: str) -> str:
    try:
        import openpyxl
    except ImportError:
        raise ImportError("XLSX support requires openpyxl: pip install openpyxl")

    try:
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    except Exception as e:
        raise ValueError(f"cannot open XLSX '{filepath}': {e}")

    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            header = f"## Sheet: {sheet_name}"
            # batch rows in groups of 50 for paragraph boundaries
            batched = []
            for i in range(0, len(rows), 50):
                batched.append("\n".join(rows[i : i + 50]))
            parts.append(f"{header}\n\n" + "\n\n".join(batched))

    wb.close()
    return "\n\n".join(parts)


def extract_csv(filepath: str) -> str:
    try:
        with open(filepath, newline="", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = []
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows.append(" | ".join(row))
    except Exception as e:
        raise ValueError(f"cannot read CSV '{filepath}': {e}")

    if not rows:
        return ""

    # batch rows in groups of 50
    parts = []
    for i in range(0, len(rows), 50):
        parts.append("\n".join(rows[i : i + 50]))
    return "\n\n".join(parts)


def extract_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("PPTX support requires python-pptx: pip install python-pptx")

    try:
        prs = Presentation(filepath)
    except Exception as e:
        raise ValueError(f"cannot open PPTX '{filepath}': {e}")

    slides = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        if texts:
            slides.append(f"[Slide {slide_num}]\n\n" + "\n\n".join(texts))

    # also extract slide notes
    for slide_num, slide in enumerate(prs.slides, 1):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slides.append(f"[Slide {slide_num} Notes]\n\n{notes}")

    return "\n\n".join(slides)


def extract_ipynb(filepath: str) -> str:
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            nb = json.load(f)
    except Exception as e:
        raise ValueError(f"cannot read notebook '{filepath}': {e}")

    cells = nb.get("cells", [])
    parts = []

    for i, cell in enumerate(cells, 1):
        cell_type = cell.get("cell_type", "")
        source = "".join(cell.get("source", []))
        if not source.strip():
            continue

        if cell_type == "markdown":
            parts.append(source)
        elif cell_type == "code":
            # detect language from notebook metadata
            lang = nb.get("metadata", {}).get("kernelspec", {}).get("language", "python")
            parts.append(f"```{lang}\n{source}\n```")

            # include text outputs (skip binary/image outputs)
            outputs = cell.get("outputs", [])
            for out in outputs:
                if out.get("output_type") == "stream":
                    text = "".join(out.get("text", []))
                    if text.strip():
                        parts.append(f"Output:\n{text.strip()}")
                elif out.get("output_type") in ("execute_result", "display_data"):
                    data = out.get("data", {})
                    if "text/plain" in data:
                        text = "".join(data["text/plain"])
                        if text.strip():
                            parts.append(f"Output:\n{text.strip()}")
        elif cell_type == "raw":
            parts.append(source)

    return "\n\n".join(parts)


def extract_code(filepath: str) -> str:
    ext = Path(filepath).suffix.lower()
    lang = LANGUAGE_MAP.get(ext, "code")

    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        code = f.read()

    # wrap in language-tagged code block for context
    return f"[Language: {lang}]\n\n```{lang}\n{code}\n```"


def extract_html(filepath: str) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("HTML support requires beautifulsoup4: pip install beautifulsoup4")

    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        raw = f.read()

    soup = BeautifulSoup(raw, "html.parser")

    # remove script and style elements
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    # extract title if present
    parts = []
    title = soup.find("title")
    if title and title.string:
        parts.append(f"# {title.string.strip()}")

    # get text with basic structure preservation
    text = soup.get_text(separator="\n")
    parts.append(text)

    return "\n\n".join(parts)


def extract_image(filepath: str) -> str:
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        raise ImportError(
            "Image OCR requires pytesseract and Pillow: "
            "pip install pytesseract pillow && sudo apt install tesseract-ocr"
        )

    try:
        pytesseract.get_tesseract_version()
    except Exception:
        raise RuntimeError(
            "Tesseract OCR binary not found. Install: sudo apt install tesseract-ocr"
        )

    img = Image.open(filepath)
    if img.mode not in ("L", "RGB"):
        img = img.convert("RGB")

    text = pytesseract.image_to_string(img, lang="eng")
    img.close()
    return text


def extract_plaintext(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


# ── dispatcher ───────────────────────────────────────────────────────

EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".csv": extract_csv,
    ".pptx": extract_pptx,
    ".png": extract_image,
    ".jpg": extract_image,
    ".jpeg": extract_image,
    ".txt": extract_plaintext,
    ".md": extract_plaintext,
    ".ipynb": extract_ipynb,
    ".html": extract_html,
    ".htm": extract_html,
    # code files
    ".py": extract_code,
    ".r": extract_code,
    ".js": extract_code,
    ".ts": extract_code,
    ".cpp": extract_code,
    ".c": extract_code,
    ".h": extract_code,
    ".hpp": extract_code,
    ".java": extract_code,
    ".sql": extract_code,
    ".go": extract_code,
    ".rs": extract_code,
    ".rb": extract_code,
    ".sh": extract_code,
    ".bash": extract_code,
}


def process_document(filepath: str) -> tuple:
    """extract text from a document file.

    returns (filename, content) tuple where content is clean UTF-8 text
    with \\n\\n paragraph separators, ready for OE ingestion.
    """
    path = Path(filepath)

    if not path.exists():
        raise FileNotFoundError(f"file not found: {filepath}")
    if not path.is_file():
        raise ValueError(f"not a file: {filepath}")

    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"unsupported format '{ext}'. supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    raw_text = EXTRACTORS[ext](str(path))
    content = normalize_paragraphs(raw_text)

    if not content.strip():
        raise ValueError(f"no text content extracted from {filepath}")

    if len(content) > MAX_CONTENT_SIZE:
        content = content[:MAX_CONTENT_SIZE]
        content += "\n\n[Document truncated at 10MB text limit]"

    return (path.name, content)


def process_document_full(filepath: str) -> dict:
    """extract text + metadata from a document file.

    returns dict with filename, content, category, summary, size_bytes, format.
    """
    path = Path(filepath)

    if not path.exists():
        raise FileNotFoundError(f"file not found: {filepath}")
    if not path.is_file():
        raise ValueError(f"not a file: {filepath}")

    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"unsupported format '{ext}'. supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    raw_text = EXTRACTORS[ext](str(path))
    content = normalize_paragraphs(raw_text)

    if not content.strip():
        raise ValueError(f"no text content extracted from {filepath}")

    if len(content) > MAX_CONTENT_SIZE:
        content = content[:MAX_CONTENT_SIZE]
        content += "\n\n[Document truncated at 10MB text limit]"

    category = get_category(ext)
    summary = generate_summary(content, path.name, category)

    return {
        "filename": path.name,
        "content": content,
        "category": category,
        "summary": summary,
        "size_bytes": path.stat().st_size,
        "format": ext.lstrip("."),
    }


# ── original file preservation ───────────────────────────────────────

def preserve_original(filepath: str, corpus_dir: str) -> dict:
    """copy original file to corpus/originals/{date}/ and create catalog entry.

    returns the catalog entry dict (without chunk_ids — caller adds those).
    """
    path = Path(filepath)
    corpus = Path(corpus_dir)
    originals_dir = corpus / "originals"

    # date-bucketed subdirectory
    date_str = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    dest_dir = originals_dir / date_str
    dest_dir.mkdir(parents=True, exist_ok=True)

    # handle filename collisions
    dest_path = dest_dir / path.name
    if dest_path.exists():
        stem = path.stem
        suffix = path.suffix
        counter = 1
        while dest_path.exists():
            dest_path = dest_dir / f"{stem}_{counter}{suffix}"
            counter += 1

    # copy file
    shutil.copy2(str(path), str(dest_path))

    ext = path.suffix.lower()
    file_id = generate_file_id()

    entry = {
        "file_id": file_id,
        "original_name": path.name,
        "stored_path": str(dest_path.relative_to(corpus)),
        "format": ext.lstrip("."),
        "category": get_category(ext),
        "size_bytes": path.stat().st_size,
        "ingested_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        "chunk_ids": [],
        "tags": [],
    }

    return entry


def update_originals_catalog(corpus_dir: str, entry: dict):
    """append entry to corpus/originals.json catalog."""
    catalog_path = Path(corpus_dir) / "originals.json"

    if catalog_path.exists():
        with open(catalog_path, "r") as f:
            catalog = json.load(f)
    else:
        catalog = []

    catalog.append(entry)

    with open(catalog_path, "w") as f:
        json.dump(catalog, f, indent=2)


def load_originals_catalog(corpus_dir: str) -> list:
    """load the originals catalog."""
    catalog_path = Path(corpus_dir) / "originals.json"
    if catalog_path.exists():
        with open(catalog_path, "r") as f:
            return json.load(f)
    return []


def get_original_by_id(corpus_dir: str, file_id: str) -> dict | None:
    """find an original file entry by its file_id."""
    catalog = load_originals_catalog(corpus_dir)
    for entry in catalog:
        if entry.get("file_id") == file_id:
            return entry
    return None


# ── server communication ─────────────────────────────────────────────

def send_to_oe(filename: str, content: str, server_url: str = OE_BASE) -> dict:
    """send extracted text to OE server for ingestion."""
    try:
        r = requests.post(
            f"{server_url}/add-file",
            json={"filename": filename, "content": content},
            timeout=120,
        )
        return r.json()
    except requests.ConnectionError:
        return {"error": f"cannot connect to OE server at {server_url}"}
    except requests.Timeout:
        return {"error": "request timed out (file may be very large)"}
    except Exception as e:
        return {"error": str(e)}


# ── CLI ───────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Ocean Eterna Document Preprocessor — convert documents to indexed text"
    )
    parser.add_argument("input", nargs="+", help="file path(s) to process")
    parser.add_argument(
        "--server",
        default=os.environ.get("OE_BASE_URL", "http://localhost:9090"),
        help="OE server URL (default: http://localhost:9090, env: OE_BASE_URL)",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="extract text and print to stdout, don't send to server",
    )
    parser.add_argument(
        "--output",
        help="write extracted text to file instead of sending to server",
    )
    args = parser.parse_args()

    success = 0
    errors = 0

    for filepath in args.input:
        print(f"\nProcessing: {filepath}")
        try:
            filename, content = process_document(filepath)
            chars = len(content)
            paragraphs = content.count("\n\n") + 1
            print(f"  extracted: {chars:,} chars, ~{paragraphs} paragraphs")

            if args.dry_run:
                print(f"\n--- {filename} ---")
                print(content[:2000])
                if len(content) > 2000:
                    print(f"\n... [{chars - 2000:,} more chars] ...")
                success += 1

            elif args.output:
                with open(args.output, "a", encoding="utf-8") as f:
                    f.write(f"\n\n=== {filename} ===\n\n{content}")
                print(f"  written to: {args.output}")
                success += 1

            else:
                result = send_to_oe(filename, content, args.server)
                if result.get("success"):
                    chunks = result.get("chunks_added", 0)
                    tokens = result.get("tokens_added", 0)
                    print(f"  ingested: {chunks} chunks, {tokens:,} tokens")
                    success += 1
                else:
                    print(f"  ERROR: {result.get('error', 'unknown error')}")
                    errors += 1

        except Exception as e:
            print(f"  ERROR: {e}")
            errors += 1

    print(f"\nDone: {success} succeeded, {errors} failed")
    if errors > 0:
        sys.exit(1)


if __name__ == "__main__":
    main()
